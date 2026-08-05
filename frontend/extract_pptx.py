from pptx import Presentation
import sys

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text_content = []
    
    for i, slide in enumerate(prs.slides):
        text_content.append(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)
    
    return "\n".join(text_content)

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract_pptx.py <path_to_pptx>")
        sys.exit(1)
    
    try:
        text = extract_text_from_pptx(sys.argv[1])
        print(text)
    except Exception as e:
        print(f"Error: {e}")
